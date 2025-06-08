from flask import Flask, request, render_template
from werkzeug.utils import secure_filename
import os

# File parsers
import docx
import PyPDF2
from pptx import Presentation

app = Flask(__name__)
UPLOAD_FOLDER = "uploads"
os.makedirs(UPLOAD_FOLDER, exist_ok=True)

@app.route("/", methods=["GET", "POST"])
def index():
    content = ""
    if request.method == "POST":
        f = request.files["file"]
        filename = secure_filename(f.filename)
        path = os.path.join(UPLOAD_FOLDER, filename)
        f.save(path)

        if filename.endswith(".pdf"):
            with open(path, "rb") as pdf:
                reader = PyPDF2.PdfReader(pdf)
                content = "\n".join([page.extract_text() or "" for page in reader.pages])
        elif filename.endswith(".docx"):
            doc = docx.Document(path)
            content = "\n".join([p.text for p in doc.paragraphs])
        elif filename.endswith(".pptx"):
            ppt = Presentation(path)
            content = "\n".join([
                shape.text for slide in ppt.slides for shape in slide.shapes if hasattr(shape, "text")
            ])
        else:
            content = "Unsupported file type."

    return render_template("index.html", content=content)

if __name__ == "__main__":
    port = int(os.environ.get("PORT", 5000))
    app.run(host="0.0.0.0", port=port)
