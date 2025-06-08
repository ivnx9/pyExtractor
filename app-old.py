from flask import Flask, request, render_template, jsonify
from werkzeug.utils import secure_filename
from flask_cors import CORS
import os

# File parsers
import docx
import PyPDF2
from pptx import Presentation

app = Flask(__name__)
CORS(app)

UPLOAD_FOLDER = "uploads"
os.makedirs(UPLOAD_FOLDER, exist_ok=True)

# Web Interface
@app.route("/", methods=["GET", "POST"])
def index():
    content = ""
    if request.method == "POST":
        f = request.files["file"]
        filename = secure_filename(f.filename)
        path = os.path.join(UPLOAD_FOLDER, filename)
        f.save(path)

        content = extract_file_content(path, filename)

    return render_template("index.html", content=content)

# API Endpoint
@app.route("/extract", methods=["POST"])
def extract():
    if "file" not in request.files:
        return jsonify({"error": "No file uploaded"}), 400

    f = request.files["file"]
    filename = secure_filename(f.filename)
    path = os.path.join(UPLOAD_FOLDER, filename)
    f.save(path)

    content = extract_file_content(path, filename)
    if content == "Unsupported file type.":
        return jsonify({"error": content}), 400

    return jsonify({"content": content})

# Core Extractor
def extract_file_content(path, filename):
    if filename.endswith(".pdf"):
        with open(path, "rb") as pdf:
            reader = PyPDF2.PdfReader(pdf)
            return "\n".join([page.extract_text() or "" for page in reader.pages])
    elif filename.endswith(".docx"):
        doc = docx.Document(path)
        return "\n".join([p.text for p in doc.paragraphs])
    elif filename.endswith(".pptx"):
        ppt = Presentation(path)
        return "\n".join([
            shape.text for slide in ppt.slides for shape in slide.shapes if hasattr(shape, "text")
        ])
    else:
        return "Unsupported file type."

if __name__ == "__main__":
    port = int(os.environ.get("PORT", 5000))
    app.run(host="0.0.0.0", port=port)
